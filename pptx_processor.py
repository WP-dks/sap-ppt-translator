from pptx import Presentation
from pptx.util import Pt
import os
import traceback

class PPTXProcessor:
    def __init__(self, translator):
        self.translator = translator
        self.korean_font = "Malgun Gothic"

    def process_presentation(self, input_path, output_path, progress_callback=None):
        try:
            from concurrent.futures import ThreadPoolExecutor
            prs = Presentation(input_path)
            
            # Step 1: Collect all text targets
            text_frames = []
            def collect_frames(shapes, is_master=False):
                if not shapes: return
                for shape in shapes:
                    # For Masters/Layouts, only translate if it's NOT a placeholder
                    # and has actual content (placeholders often overlap or cause shifts)
                    if is_master and (shape.is_placeholder or not shape.has_text_frame):
                        continue
                        
                    if shape.has_text_frame:
                        text_frames.append(shape.text_frame)
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                text_frames.append(cell.text_frame)
                    if shape.shape_type == 6:  # Group shape
                        collect_frames(shape.shapes, is_master)
            
            # 1.1: Collect from slides (Higher priority)
            for slide in prs.slides:
                collect_frames(slide.shapes, is_master=False)
                if slide.has_notes_slide:
                    collect_frames(slide.notes_slide.shapes, is_master=False)
            
            # 1.2: Collect from layouts/masters (Only if user wanted 'thorough' but stay safe)
            for master in prs.slide_masters:
                collect_frames(master.shapes, is_master=True)
                for layout in master.slide_layouts:
                    collect_frames(layout.shapes, is_master=True)
            
            # Step 2: Extract all unique paragraphs
            unique_texts = set()
            paragraphs_to_translate = []
            for tf in text_frames:
                for p in tf.paragraphs:
                    if p.text.strip() and len(p.text.strip()) > 1: # Skip single chars/bullets
                        unique_texts.add(p.text)
                        paragraphs_to_translate.append(p)

            # Step 3: Translate unique texts in parallel
            total_unique = len(unique_texts)
            translation_map = {}
            translation_errors = []
            processed_count = 0

            # Only translate if there's actually something to translate
            if unique_texts:
                max_workers = 15
                with ThreadPoolExecutor(max_workers=max_workers) as executor:
                    future_to_text = {executor.submit(self.translator.translate, text): text for text in unique_texts}
                    
                    for i, future in enumerate(future_to_text):
                        text = future_to_text[future]
                        try:
                            result = future.result()
                            if result:
                                translation_map[text] = result
                            else:
                                translation_map[text] = text
                        except Exception as e:
                            err_msg = f"'{text[:30]}...' 번역 실패 ({type(e).__name__}): {str(e)}"
                            translation_errors.append(err_msg)
                            translation_map[text] = text
                        
                        processed_count += 1
                        if progress_callback:
                            progress_callback((processed_count / total_unique) * 0.8)

            # Step 4: Apply translations
            for i, p in enumerate(paragraphs_to_translate):
                if p.text in translation_map:
                    translated_text = translation_map[p.text]
                    if translated_text and translated_text != p.text:
                        self._update_paragraph_text(p, translated_text)
                
                if progress_callback:
                    progress_callback(0.8 + (i / len(paragraphs_to_translate)) * 0.2)

            prs.save(output_path)
            return output_path, translation_errors
        except Exception as e:
            print(f"Critical error in process_presentation: {e}")
            traceback.print_exc()
            raise e

    def _update_paragraph_text(self, paragraph, translated_text):
        if translated_text is None:
            return
        
        try:
            # Capture properties from the runs before clearing
            font_size = None
            font_bold = None
            font_italic = None
            font_color_rgb = None
            font_color_theme = None
            
            # Capture critical paragraph properties to prevent shifts
            para_alignment = paragraph.alignment
            para_level = paragraph.level
            para_line_spacing = paragraph.line_spacing
            para_space_before = paragraph.space_before
            para_space_after = paragraph.space_after

            for run in paragraph.runs:
                f = run.font
                if font_size is None and f.size: font_size = f.size
                if font_bold is None and f.bold is not None: font_bold = f.bold
                if font_italic is None and f.italic is not None: font_italic = f.italic
                
                try:
                    if font_color_rgb is None and f.color.type == 1:
                        font_color_rgb = f.color.rgb
                    elif font_color_theme is None and f.color.type == 2:
                        font_color_theme = f.color.theme_color
                except: pass

            paragraph.text = translated_text
            
            # Apply paragraph properties back exactly
            paragraph.alignment = para_alignment
            paragraph.level = para_level
            if para_line_spacing is not None: paragraph.line_spacing = para_line_spacing
            if para_space_before is not None: paragraph.space_before = para_space_before
            if para_space_after is not None: paragraph.space_after = para_space_after
            
            # Apply run properties
            if paragraph.runs:
                run = paragraph.runs[0]
                run.font.name = self.korean_font
                if font_size: run.font.size = font_size
                if font_bold is not None: run.font.bold = font_bold
                if font_italic is not None: run.font.italic = font_italic
                
                if font_color_rgb is not None:
                    run.font.color.rgb = font_color_rgb
                elif font_color_theme is not None:
                    run.font.color.theme_color = font_color_theme
        except Exception as e:
            print(f"Error updating paragraph text: {e}")
            if translated_text is not None:
                paragraph.text = translated_text
